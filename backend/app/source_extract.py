"""
Server-side source text extraction for uploaded files.
Honest errors when a parser is unavailable — never fake content.
"""

from __future__ import annotations

import io
from typing import Any


MAX_EXTRACT_CHARS = 45_000
MAX_UPLOAD_BYTES = 10 * 1024 * 1024  # 10 MB


def _clip(text: str) -> str:
    text = (text or "").strip()
    if len(text) > MAX_EXTRACT_CHARS:
        return text[:MAX_EXTRACT_CHARS]
    return text


def extract_text_from_bytes(
    *,
    filename: str,
    content_type: str,
    data: bytes,
) -> dict[str, Any]:
    if not data:
        return {"success": False, "error": "Empty file."}
    if len(data) > MAX_UPLOAD_BYTES:
        return {"success": False, "error": "File exceeds the 10MB limit."}

    name = (filename or "upload").lower()
    ctype = (content_type or "").lower()

    # Plain text family
    if (
        name.endswith((".txt", ".md", ".markdown", ".csv", ".json", ".log"))
        or ctype.startswith("text/")
        or ctype in {"application/json", "text/csv", "text/markdown"}
    ):
        try:
            text = data.decode("utf-8-sig")  # strip BOM if present
        except UnicodeDecodeError:
            text = data.decode("latin-1", errors="replace")
        text = _clip(text)
        if not text:
            return {"success": False, "error": "No extractable text found in this file."}
        return {"success": True, "text": text, "format": "text"}

    # PDF
    if name.endswith(".pdf") or "pdf" in ctype:
        try:
            from pypdf import PdfReader  # type: ignore
        except ImportError:
            return {
                "success": False,
                "error": "PDF extraction is not available on this server (pypdf missing).",
            }
        try:
            reader = PdfReader(io.BytesIO(data))
            parts: list[str] = []
            for page in reader.pages:
                parts.append(page.extract_text() or "")
            text = _clip("\n".join(parts))
            if not text:
                return {
                    "success": False,
                    "error": "PDF contained no extractable text (it may be image-only). OCR is not configured.",
                }
            return {"success": True, "text": text, "format": "pdf", "pages": len(reader.pages)}
        except Exception as exc:  # noqa: BLE001
            return {"success": False, "error": f"Failed to read PDF: {str(exc)[:200]}"}

    # DOCX
    if name.endswith(".docx") or "wordprocessingml" in ctype:
        try:
            import docx  # type: ignore
        except ImportError:
            return {
                "success": False,
                "error": "DOCX extraction is not available on this server (python-docx missing).",
            }
        try:
            document = docx.Document(io.BytesIO(data))
            text = _clip("\n".join(p.text for p in document.paragraphs if p.text))
            if not text:
                return {"success": False, "error": "DOCX contained no extractable text."}
            return {"success": True, "text": text, "format": "docx"}
        except Exception as exc:  # noqa: BLE001
            return {"success": False, "error": f"Failed to read DOCX: {str(exc)[:200]}"}

    # XLSX
    if name.endswith((".xlsx", ".xlsm")) or "spreadsheetml" in ctype:
        try:
            from openpyxl import load_workbook  # type: ignore
        except ImportError:
            return {
                "success": False,
                "error": "Spreadsheet extraction is not available on this server (openpyxl missing).",
            }
        try:
            wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
            lines: list[str] = []
            for sheet in wb.worksheets:
                lines.append(f"# Sheet: {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    cells = ["" if c is None else str(c) for c in row]
                    if any(cells):
                        lines.append("\t".join(cells))
            text = _clip("\n".join(lines))
            if not text:
                return {"success": False, "error": "Spreadsheet contained no extractable text."}
            return {"success": True, "text": text, "format": "xlsx"}
        except Exception as exc:  # noqa: BLE001
            return {"success": False, "error": f"Failed to read spreadsheet: {str(exc)[:200]}"}

    # PPTX
    if name.endswith(".pptx") or "presentationml" in ctype:
        try:
            from pptx import Presentation  # type: ignore
        except ImportError:
            return {
                "success": False,
                "error": "PPTX extraction is not available on this server (python-pptx missing).",
            }
        try:
            prs = Presentation(io.BytesIO(data))
            parts: list[str] = []
            for i, slide in enumerate(prs.slides, start=1):
                parts.append(f"# Slide {i}")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        parts.append(shape.text)
            text = _clip("\n".join(parts))
            if not text:
                return {"success": False, "error": "Presentation contained no extractable text."}
            return {"success": True, "text": text, "format": "pptx"}
        except Exception as exc:  # noqa: BLE001
            return {"success": False, "error": f"Failed to read PPTX: {str(exc)[:200]}"}

    # Images / audio — require services not configured by default
    if name.endswith((".png", ".jpg", ".jpeg", ".webp", ".gif")) or ctype.startswith("image/"):
        return {
            "success": False,
            "error": "Image OCR is not configured. Paste text from the image or upload a text/PDF document.",
        }
    if name.endswith((".mp3", ".wav", ".m4a", ".ogg")) or ctype.startswith("audio/"):
        return {
            "success": False,
            "error": "Audio transcription is not configured. Provide a transcript or text notes instead.",
        }
    if name.endswith((".ppt", ".doc", ".xls")):
        return {
            "success": False,
            "error": "Legacy Office binary formats are not supported. Convert to DOCX, PPTX, XLSX, PDF, or TXT.",
        }

    return {
        "success": False,
        "error": f"Unsupported file type for extraction: {filename or content_type or 'unknown'}.",
    }
